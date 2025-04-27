import streamlit as st
import os
import tempfile
import time
import asyncio
import aiohttp
import json
import hashlib
import numpy as np
from typing import List, Dict, Any, Tuple, Optional
import pandas as pd

# Document processing libraries
import PyPDF2
import docx
import pptx
from pptx import Presentation
import openpyxl
import re

# Set page configuration
st.set_page_config(
    page_title="Notesight",
    page_icon="📚",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Cache directory for processed chunks
CACHE_DIR = "./.cache"
os.makedirs(CACHE_DIR, exist_ok=True)

class DocumentProcessor:
    """Handles document processing for various file types"""
    
    @staticmethod
    def extract_text_from_pdf(file_path: str) -> List[str]:
        """Extract text from PDF files with page-wise extraction"""
        text_chunks = []
        
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            total_pages = len(pdf_reader.pages)
            
            # Extract larger chunks (20 pages at a time)
            for i in range(0, total_pages, 20):
                batch_text = ""
                end_page = min(i + 20, total_pages)
                
                for page_num in range(i, end_page):
                    page = pdf_reader.pages[page_num]
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            batch_text += f"\n\n--- Page {page_num + 1} ---\n\n{page_text}"
                    except Exception as e:
                        batch_text += f"\n\n--- Page {page_num + 1} [Error: {str(e)}] ---\n\n"
                
                if batch_text.strip():
                    text_chunks.append(batch_text)
        
        return text_chunks

    @staticmethod
    def extract_text_from_docx(file_path: str) -> List[str]:
        """Extract text from DOCX files preserving structure"""
        doc = docx.Document(file_path)
        text_chunks = []
        chunk = ""
        chunk_size = 15000  # Increased chunk size for fewer API calls
        
        for para in doc.paragraphs:
            para_text = para.text.strip()
            if para_text:
                if len(chunk) + len(para_text) > chunk_size and chunk:
                    text_chunks.append(chunk)
                    chunk = ""
                
                chunk += para_text + "\n\n"
        
        if chunk:
            text_chunks.append(chunk)
            
        return text_chunks

    @staticmethod
    def extract_text_from_pptx(file_path: str) -> List[str]:
        """Extract text from PPTX presentations"""
        presentation = Presentation(file_path)
        all_text = ""
        slide_count = 0
        
        for i, slide in enumerate(presentation.slides):
            slide_text = f"--- Slide {i+1} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text += shape.text + "\n"
            
            if len(slide_text) > 20:  # More than just the slide header
                all_text += slide_text + "\n\n"
                slide_count += 1
                
            # Create a chunk every 10 slides or at the end
            if slide_count % 10 == 0 or i == len(presentation.slides) - 1:
                if all_text.strip():
                    return [all_text]
                all_text = ""
        
        return [all_text] if all_text.strip() else []

    @staticmethod
    def extract_text_from_xlsx(file_path: str) -> List[str]:
        """Extract text from Excel files"""
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        all_sheets_text = ""
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_text = f"--- Sheet: {sheet_name} ---\n\n"
            
            # Get the dimensions of the sheet
            rows = list(sheet.rows)
            if not rows:
                continue
                
            col_names = [str(cell.value) if cell.value is not None else "" for cell in rows[0]]
            
            # Process sheet content
            data_rows = []
            for row in rows[1:]:
                cell_values = [str(cell.value) if cell.value is not None else "" for cell in row]
                if any(cell.strip() for cell in cell_values):  # Skip empty rows
                    data_rows.append(cell_values)
            
            # Convert to tabular format for better context
            try:
                df = pd.DataFrame(data_rows, columns=col_names)
                sheet_text += df.to_string(index=False) + "\n\n"
                all_sheets_text += sheet_text
            except Exception as e:
                # Fallback to manual formatting if pandas conversion fails
                sheet_text += "| " + " | ".join(col_names) + " |\n"
                sheet_text += "| " + " | ".join(["---" for _ in col_names]) + " |\n"
                
                for row in data_rows:
                    sheet_text += "| " + " | ".join(row) + " |\n"
                
                all_sheets_text += sheet_text
        
        return [all_sheets_text] if all_sheets_text.strip() else []


class OpenAIClient:
    """Async client for OpenAI API calls"""
    
    def __init__(self, api_key: str):
        self.api_key = api_key
        self.headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}"
        }
        self.embedding_cache = {}
        self.semaphore = asyncio.Semaphore(5)  # Limit concurrent API calls
        
    async def get_embedding(self, session: aiohttp.ClientSession, text: str) -> List[float]:
        """Get embedding vector for text using OpenAI's text-embedding-3-small model"""
        # Check cache first
        cache_key = hashlib.md5(text.encode()).hexdigest()
        if cache_key in self.embedding_cache:
            return self.embedding_cache[cache_key]
            
        # Clean and truncate the text if necessary
        text = text.strip()
        if len(text) > 8000:  # Embedding model has a token limit
            text = text[:8000]
            
        async with self.semaphore:
            url = "https://api.openai.com/v1/embeddings"
            payload = {
                "model": "text-embedding-3-small",  # Use smaller model for speed
                "input": text
            }
            
            async with session.post(url, json=payload, headers=self.headers) as response:
                if response.status == 200:
                    data = await response.json()
                    embedding = data["data"][0]["embedding"]
                    # Cache the result
                    self.embedding_cache[cache_key] = embedding
                    return embedding
                else:
                    error_text = await response.text()
                    raise Exception(f"Error getting embedding: {error_text}")
    
    async def generate_completions(
        self,
        session: aiohttp.ClientSession,
        messages: List[Dict[str, str]],
        model: str = "gpt-3.5-turbo",  # Changed default to 3.5-turbo
        temperature: float = 0.7,
        max_tokens: int = 1800
    ) -> str:
        """Generate text completions using OpenAI's Chat Completion API"""
        async with self.semaphore:
            url = "https://api.openai.com/v1/chat/completions"
            payload = {
                "model": model,
                "messages": messages,
                "temperature": temperature,
                "max_tokens": max_tokens
            }
            
            async with session.post(url, json=payload, headers=self.headers) as response:
                if response.status == 200:
                    data = await response.json()
                    return data["choices"][0]["message"]["content"]
                else:
                    error_text = await response.text()
                    raise Exception(f"Error generating completion: {error_text}")


class FastChunker:
    """Creates chunks quickly without relying on embeddings"""
    
    @staticmethod
    def create_chunks(text_chunks: List[str], max_chunk_size: int = 12000) -> List[str]:
        """Create chunks without embeddings for faster processing"""
        all_chunks = []
        
        for text in text_chunks:
            # If chunk is already small enough, keep it as is
            if len(text) <= max_chunk_size:
                all_chunks.append(text)
                continue
            
            # Split into paragraphs for larger texts
            paragraphs = text.split("\n\n")
            
            current_chunk = ""
            for para in paragraphs:
                if len(current_chunk) + len(para) > max_chunk_size and current_chunk:
                    all_chunks.append(current_chunk)
                    current_chunk = para
                else:
                    current_chunk += "\n\n" + para if current_chunk else para
            
            if current_chunk:
                all_chunks.append(current_chunk)
        
        return all_chunks


class QAGenerator:
    """Generates MCQs and flashcards using OpenAI API"""
    
    def __init__(self, client: OpenAIClient):
        self.client = client
        
    def generate_chunk_hash(self, chunk: str) -> str:
        """Generate a hash for a text chunk for caching"""
        return hashlib.md5(chunk.encode()).hexdigest()
    
    def check_cache(self, chunk_hash: str, qa_type: str) -> Optional[Dict]:
        """Check if we have cached results for this chunk"""
        cache_file = os.path.join(CACHE_DIR, f"{chunk_hash}_{qa_type}.json")
        if os.path.exists(cache_file):
            with open(cache_file, 'r') as f:
                return json.load(f)
        return None
    
    def save_to_cache(self, chunk_hash: str, qa_type: str, result: Dict) -> None:
        """Save results to cache"""
        cache_file = os.path.join(CACHE_DIR, f"{chunk_hash}_{qa_type}.json")
        with open(cache_file, 'w') as f:
            json.dump(result, f)
    
    async def generate_mcqs(
        self, 
        session: aiohttp.ClientSession,
        chunk: str, 
        num_questions: int = 5
    ) -> Dict:
        """Generate MCQs from a text chunk"""
        chunk_hash = self.generate_chunk_hash(chunk)
        
        # Check cache first
        cached_result = self.check_cache(chunk_hash, "mcq")
        if cached_result:
            return cached_result
        
        prompt = f"""
        I need you to generate {num_questions} multiple-choice questions based on the following text. 
        For each question:
        1. Create a clear, specific question that tests understanding of key concepts
        2. Provide 4 options (A, B, C, D) with only one correct answer
        3. Mark the correct answer
        4. Add a brief explanation of why the marked answer is correct
        
        Focus on important concepts and avoid trivial details. Ensure questions test understanding rather than mere recall.
        
        TEXT:
        {chunk}
        
        FORMAT YOUR RESPONSE AS JSON:
        {{
            "questions": [
                {{
                    "question": "Question text",
                    "options": ["A. Option 1", "B. Option 2", "C. Option 3", "D. Option 4"],
                    "correct_answer": "A",
                    "explanation": "Explanation text"
                }},
                ...
            ]
        }}
        """
        
        try:
            messages = [
                {"role": "system", "content": "You are an expert educator who creates high-quality assessment questions. Respond with valid JSON only."},
                {"role": "user", "content": prompt}
            ]
            
            result_text = await self.client.generate_completions(
                session=session,
                messages=messages,
                model="gpt-3.5-turbo-1106",  # Specific model for JSON output
                temperature=0.5,
                max_tokens=2000
            )
            
            # Extract JSON from the response
            json_match = re.search(r'```json\s*(.*?)\s*```', result_text, re.DOTALL)
            if json_match:
                result_text = json_match.group(1)
            
            try:
                result = json.loads(result_text)
                self.save_to_cache(chunk_hash, "mcq", result)
                return result
            except json.JSONDecodeError:
                # Fallback if JSON is malformed
                return {"questions": [], "error": "Failed to parse response as JSON"}
                
        except Exception as e:
            return {"questions": [], "error": str(e)}
    
    async def generate_flashcards(
        self, 
        session: aiohttp.ClientSession,
        chunk: str, 
        num_cards: int = 8
    ) -> Dict:
        """Generate flashcards from a text chunk"""
        chunk_hash = self.generate_chunk_hash(chunk)
        
        # Check cache first
        cached_result = self.check_cache(chunk_hash, "flashcard")
        if cached_result:
            return cached_result
        
        prompt = f"""
        Create {num_cards} flashcards based on the following text. Each flashcard should:
        1. Focus on a key concept, definition, or important fact
        2. Have a clear front (question/prompt) and back (answer/explanation)
        3. Be concise but comprehensive
        
        Prioritize the most important information from the text.
        
        TEXT:
        {chunk}
        
        FORMAT YOUR RESPONSE AS JSON:
        {{
            "flashcards": [
                {{
                    "front": "concept or topic",
                    "back": "Answer or explanation"
                }},
                ...
            ]
        }}
        """
        
        try:
            messages = [
                {"role": "system", "content": "You are an expert educator who creates clear, effective study materials. Respond with valid JSON only."},
                {"role": "user", "content": prompt}
            ]
            
            result_text = await self.client.generate_completions(
                session=session,
                messages=messages,
                model="gpt-3.5-turbo-1106",  # Specific model for JSON output
                temperature=0.5,
                max_tokens=2000
            )
            
            # Extract JSON from the response
            json_match = re.search(r'```json\s*(.*?)\s*```', result_text, re.DOTALL)
            if json_match:
                result_text = json_match.group(1)
            
            try:
                result = json.loads(result_text)
                self.save_to_cache(chunk_hash, "flashcard", result)
                return result
            except json.JSONDecodeError:
                # Fallback if JSON is malformed
                return {"flashcards": [], "error": "Failed to parse response as JSON"}
                
        except Exception as e:
            return {"flashcards": [], "error": str(e)}


async def process_document(
    file_path: str, 
    file_extension: str, 
    progress_bar,
    status_text
) -> List[str]:
    """Process document and extract text chunks for later use"""
    # Extract text chunks from the document
    status_text.text("Extracting text from document...")
    if file_extension == '.pdf':
        initial_chunks = DocumentProcessor.extract_text_from_pdf(file_path)
    elif file_extension == '.docx':
        initial_chunks = DocumentProcessor.extract_text_from_docx(file_path)
    elif file_extension == '.pptx':
        initial_chunks = DocumentProcessor.extract_text_from_pptx(file_path)
    elif file_extension in ['.xlsx', '.xls']:
        initial_chunks = DocumentProcessor.extract_text_from_xlsx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_extension}")
    
    # Create optimized chunks
    status_text.text("Creating document chunks...")
    optimized_chunks = FastChunker.create_chunks(initial_chunks)
    
    progress_bar.progress(100)
    status_text.text("Document processed successfully!")
    
    return optimized_chunks


async def generate_mcqs(
    chunks: List[str],
    api_key: str,
    mcq_count: int,
    progress_bar,
    status_text
) -> List[Dict]:
    """Generate MCQs from text chunks"""
    client = OpenAIClient(api_key)
    qa_generator = QAGenerator(client)
    
    progress_bar.progress(10)
    
    # Calculate how many questions to generate per chunk
    total_chunks = len(chunks)
    questions_per_chunk = max(1, mcq_count // total_chunks + 1)
    
    # Setup progress tracking
    progress_bar.progress(20)
    status_text.text("Generating MCQs chunks...")
    
    # Create tasks for parallel processing
    async with aiohttp.ClientSession() as session:
        # Create a list of tasks for each chunk
        mcq_tasks = []
        for chunk in chunks:
            mcq_tasks.append(qa_generator.generate_mcqs(session, chunk, questions_per_chunk))
        
        # Process chunks in parallel
        mcq_results = await asyncio.gather(*mcq_tasks)
        
        progress_bar.progress(90)
        
        # Combine all MCQs
        all_mcqs = []
        for result in mcq_results:
            if "questions" in result and result["questions"]:
                all_mcqs.extend(result["questions"])
        
        # Limit to requested count
        if len(all_mcqs) > mcq_count:
            all_mcqs = all_mcqs[:mcq_count]
        
        progress_bar.progress(100)
        status_text.text("Done!")
        
        return all_mcqs


async def generate_flashcards(
    chunks: List[str],
    api_key: str,
    flashcard_count: int,
    progress_bar,
    status_text
) -> List[Dict]:
    """Generate flashcards from text chunks"""
    client = OpenAIClient(api_key)
    qa_generator = QAGenerator(client)
    
    progress_bar.progress(10)
    
    # Calculate how many flashcards to generate per chunk
    total_chunks = len(chunks)
    cards_per_chunk = max(1, flashcard_count // total_chunks + 1)
    
    # Setup progress tracking
    progress_bar.progress(20)
    status_text.text("Generating flashcards...")
    
    # Create tasks for parallel processing
    async with aiohttp.ClientSession() as session:
        # Create a list of tasks for each chunk
        flashcard_tasks = []
        for chunk in chunks:
            flashcard_tasks.append(qa_generator.generate_flashcards(session, chunk, cards_per_chunk))
        
        # Process chunks in parallel
        flashcard_results = await asyncio.gather(*flashcard_tasks)
        
        progress_bar.progress(90)
        
        # Combine all flashcards
        all_flashcards = []
        for result in flashcard_results:
            if "flashcards" in result and result["flashcards"]:
                all_flashcards.extend(result["flashcards"])
        
        # Limit to requested count
        if len(all_flashcards) > flashcard_count:
            all_flashcards = all_flashcards[:flashcard_count]
        
        progress_bar.progress(100)
        status_text.text("Done!")
        
        return all_flashcards


# Streamlit app  
def main():
    st.title("📚 Notesight")
    st.markdown("""
    Upload your document (PDF, DOCX, PPTX, XLSX) 
    """)

    # Sidebar for configuration
    st.sidebar.header("Configuration")

    # OpenAI API key input
    api_key = st.sidebar.text_input("OpenAI API Key", type="password")

    # Generation options
    st.sidebar.subheader("Generation Options")
    gen_option = st.sidebar.radio(
        "What would you like to generate?",
        ["Multiple Choice Questions", "Flashcards"]
    )
    
    # Question count settings based on selection
    if gen_option == "Multiple Choice Questions":
        item_count = st.sidebar.number_input("Number of MCQs", min_value=5, max_value=50, value=10)
    else:
        item_count = st.sidebar.number_input("Number of Flashcards", min_value=5, max_value=50, value=10)

    # File uploader
    uploaded_file = st.file_uploader("Upload your document", type=['pdf', 'docx', 'pptx', 'xlsx'])

    # Initialize session state for storing processed document chunks
    if 'processed_chunks' not in st.session_state:
        st.session_state.processed_chunks = None
        
    if 'generated_mcqs' not in st.session_state:
        st.session_state.generated_mcqs = None
        
    if 'generated_flashcards' not in st.session_state:
        st.session_state.generated_flashcards = None

    if uploaded_file is not None:
        # Display file info
        file_details = {
            "Filename": uploaded_file.name,
            "File size": f"{uploaded_file.size / (1024 * 1024):.2f} MB"
        }
        st.write(file_details)
        
        # Check if API key is provided
        if not api_key:
            st.warning("Please enter your OpenAI API key in the sidebar.")
            st.stop()
        
        # Combined process & generate button
        if st.button(f"Process Document & Generate {gen_option}"):
            file_extension = os.path.splitext(uploaded_file.name)[1].lower()
            
            # Save uploaded file temporarily
            with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as tmp_file:
                tmp_file.write(uploaded_file.getvalue())
                temp_file_path = tmp_file.name
            
            # Initialize progress displays
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            try:
                # Run the async processing
                loop = asyncio.new_event_loop()
                asyncio.set_event_loop(loop)
                
                # Process document
                st.session_state.processed_chunks = loop.run_until_complete(
                    process_document(
                        temp_file_path, 
                        file_extension,
                        progress_bar,
                        status_text
                    )
                )
                
                # Generate QA content
                if gen_option == "Multiple Choice Questions":
                    st.session_state.generated_mcqs = loop.run_until_complete(
                        generate_mcqs(
                            st.session_state.processed_chunks,
                            api_key,
                            item_count,
                            progress_bar,
                            status_text
                        )
                    )
                else:
                    st.session_state.generated_flashcards = loop.run_until_complete(
                        generate_flashcards(
                            st.session_state.processed_chunks,
                            api_key,
                            item_count,
                            progress_bar,
                            status_text
                        )
                    )
                
                st.success(f"{gen_option} generated successfully!")
                
            except Exception as e:
                st.error(f"Error: {str(e)}")
            
            finally:
                # Clean up the temporary file
                try:
                    os.unlink(temp_file_path)
                except:
                    pass

        # Display results if available
        if st.session_state.generated_mcqs is not None and gen_option == "Multiple Choice Questions":
            st.subheader(f"Generated MCQs ({len(st.session_state.generated_mcqs)})")
            
            for i, mcq in enumerate(st.session_state.generated_mcqs, 1):
                with st.expander(f"Question {i}: {mcq['question'][:100]}..."):
                    st.write(f"**{mcq['question']}**")
                    
                    for option in mcq['options']:
                        if option.startswith(mcq['correct_answer'] + "."):
                            st.success(option)
                        else:
                            st.write(option)
                    
                    st.info(f"**Explanation:** {mcq['explanation']}")
            
            # MCQ Export
            mcq_json = json.dumps({"mcqs": st.session_state.generated_mcqs}, indent=2)
            st.download_button(
                label="Download MCQs (JSON)",
                data=mcq_json,
                file_name=f"{uploaded_file.name}_mcqs.json",
                mime="application/json"
            )
                
        elif st.session_state.generated_flashcards is not None and gen_option == "Flashcards":
            st.subheader(f"Generated Flashcards ({len(st.session_state.generated_flashcards)})")
            
            for i, card in enumerate(st.session_state.generated_flashcards, 1):
                with st.expander(f"Card {i}: {card['front'][:100]}..."):
                    st.write(f"**Front:** {card['front']}")
                    st.success(f"**Back:** {card['back']}")
            
            # Flashcard Export
            flashcard_json = json.dumps({"flashcards": st.session_state.generated_flashcards}, indent=2)
            st.download_button(
                label="Download Flashcards (JSON)",
                data=flashcard_json,
                file_name=f"{uploaded_file.name}_flashcards.json",
                mime="application/json"
            )

    # Footer
    st.markdown("---")

# Run the app
if __name__ == "__main__":
    main()